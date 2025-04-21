import io
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import re

def extract_text_from_shape(shape):
    """Extract text from any PowerPoint shape including tables and groups"""
    text = ""
    
    # Extract text from shape if it has text
    if hasattr(shape, "text") and shape.text.strip():
        text += shape.text.strip() + "\n"
    
    # Extract text from tables
    if shape.has_table:
        table = shape.table
        rows = []
        for r_idx in range(0, len(table.rows)):
            row = []
            for c_idx in range(0, len(table.columns)):
                cell = table.cell(r_idx, c_idx)
                # Add cell text with appropriate spacing
                row.append(cell.text.strip())
            rows.append(" | ".join(row))
        text += "\n".join(rows) + "\n"
    
    # Extract text from charts
    if shape.has_chart:
        chart = shape.chart
        # Get chart title if available
        if hasattr(chart, 'chart_title') and chart.chart_title:
            text += f"CHART: {chart.chart_title.text_frame.text}\n"
        
        # Extract data from chart
        if hasattr(chart, 'series'):
            series_names = []
            for series in chart.series:
                if hasattr(series, 'name'):
                    series_names.append(str(series.name))
                
                # Try to get data values
                values = []
                try:
                    if hasattr(series, 'values'):
                        for val in series.values:
                            values.append(str(val))
                    text += f"Series {', '.join(series_names)}: {', '.join(values)}\n"
                except:
                    pass
    
    # Extract text from grouped shapes
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for subshape in shape.shapes:
            text += extract_text_from_shape(subshape)
    
    return text

def identify_slide_type(slide_text):
    """Identify the type of slide based on its content"""
    slide_text_upper = slide_text.upper()
    
    if "ACCOUNT PERFORMANCE REPORT" in slide_text_upper or "PREPARED FOR" in slide_text_upper:
        return "INTRO"
    elif "SEARCH OVERVIEW" in slide_text_upper:
        return "SEARCH_OVERVIEW"
    elif "SEARCH CAMPAIGNS" in slide_text_upper:
        return "SEARCH_CAMPAIGNS"
    elif "SEARCH KEYWORDS" in slide_text_upper or "TOP SEARCH KEYWORDS" in slide_text_upper:
        return "SEARCH_KEYWORDS"
    # Be very specific about PMAX w/ VLA to differentiate it from regular PMAX
    elif "PERFORMANCEMAX W/ VLA" in slide_text_upper:
        return "PMAX_VLA"
    elif "PERFORMANCEMAX" in slide_text_upper:
        return "PMAX"
    elif "SOCIAL ADS SUMMARY" in slide_text_upper:
        return "SOCIAL_SUMMARY"
    elif "SOCIAL CAMPAIGNS" in slide_text_upper or "SOCIAL ADS" in slide_text_upper:
        return "SOCIAL_CAMPAIGNS"
    elif "VIDEO" in slide_text_upper:
        return "VIDEO"
    # More specific matching for BCDF
    elif "BCDF PROGRAM" in slide_text_upper or "BUSINESS CENTER DIRECTED FUNDS" in slide_text_upper or "STELLANTIS" in slide_text_upper:
        return "BCDF"
    elif "DEMAND GEN" in slide_text_upper:
        return "DEMAND_GEN"
    else:
        return "OTHER"

def process_slide_for_metrics(slide_text, slide_type):
    """Enhance the slide text by marking potential metrics for easier AI extraction"""
    processed_lines = []
    
    # Split content into lines
    lines = slide_text.split("\n")
    
    for line in lines:
        # Format metric lines with improved pattern matching
        if re.search(r'\d{1,3}(?:,\d{3})*\s*(?:Impressions|Clicks|Conversions|Views|Reach|VDP\s+Views)', line, re.IGNORECASE):
            processed_lines.append(f"METRIC: {line}")
        elif re.search(r'\$\d+\.\d{2}\s*(?:CPC|CPM|Cost|Cost / Conversion|Cost per Conversion)', line, re.IGNORECASE):
            processed_lines.append(f"METRIC: {line}")
        elif re.search(r'\d+\.\d+%\s*(?:CTR|Conversion Rate|View Rate)', line, re.IGNORECASE):
            processed_lines.append(f"METRIC: {line}")
        # Add pattern for campaign names and their metrics
        elif re.search(r'(New_AIA|Used\|CPO|PMax_|VLA_)\w+', line, re.IGNORECASE):
            processed_lines.append(f"CAMPAIGN: {line}")
        else:
            processed_lines.append(line)
    
    # Add additional section-specific markup
    if slide_type == "PMAX_VLA":
        # Mark VLA-specific content
        processed_lines = [f"VLA_CONTENT: {line}" if not line.startswith("VLA_CONTENT:") and "VLA_" in line 
                            else line for line in processed_lines]
    
    if slide_type == "BCDF":
        # Mark BCDF-specific content
        processed_lines = [f"BCDF_CONTENT: {line}" if not line.startswith("BCDF_CONTENT:") 
                            else line for line in processed_lines]
    
    return "\n".join(processed_lines)

def extract_text_from_pptx(uploaded_file):
    """Extract structured text from a PowerPoint file optimized for AI processing"""
    file_bytes = uploaded_file.getvalue()
    prs = Presentation(io.BytesIO(file_bytes))
    
    # Extract structured information from each slide
    structured_text_parts = []
    
    for i, slide in enumerate(prs.slides):
        # Extract all text from the slide
        slide_text = ""
        for shape in slide.shapes:
            slide_text += extract_text_from_shape(shape)
        
        # Identify slide type
        slide_type = identify_slide_type(slide_text)
        
        # Add clear section markers for the AI
        section_header = f"--- SLIDE {i+1} | TYPE: {slide_type} ---"
        structured_text_parts.append(section_header)
        
        # Process slide content for better AI recognition
        processed_text = process_slide_for_metrics(slide_text, slide_type)
        structured_text_parts.append(processed_text)
        
        # Add a separator between slides
        structured_text_parts.append("-" * 80)
    
    # Generate the complete structured text optimized for AI processing
    structured_text = "\n\n".join(structured_text_parts)
    
    return structured_text
