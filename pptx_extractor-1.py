import io, re
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# ----------------------------------------------------------------------------
# Helpers to pull raw text
# ----------------------------------------------------------------------------
def extract_text_from_shape(shape):
    txt = ""
    # plain textbox
    if hasattr(shape, "text") and shape.text.strip():
        txt += shape.text.strip() + "\n"
    # table cells
    if shape.has_table:
        rows = []
        for r in shape.table.rows:
            rows.append(" | ".join(c.text.strip() for c in r.cells))
        txt += "\n".join(rows) + "\n"
    # (optional) chart titles / series
    if shape.has_chart and shape.chart.chart_title:
        txt += f"CHART: {shape.chart.chart_title.text_frame.text}\n"
    # grouped shapes
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for sub in shape.shapes:
            txt += extract_text_from_shape(sub)
    return txt

# ----------------------------------------------------------------------------
# Slide classification
# ----------------------------------------------------------------------------
def identify_slide_type(text):
    t = text.upper()
    if "PERFORMANCEMAX W/ VLA" in t:            return "PMAX_VLA"
    if "PERFORMANCEMAX" in t:                  return "PMAX"
    if "SOCIAL ADS" in t:                      return "SOCIAL"
    if "DEMAND GEN" in t:                      return "DEMAND_GEN"
    if "VIDEO" in t and "DISPLAY" in t:        return "VIDEO"
    if "BCDF" in t or "BUSINESS CENTER DIRECTED FUNDS" in t: return "BCDF"
    return "OTHER"

# ----------------------------------------------------------------------------
# Simple regex extract helpers
# ----------------------------------------------------------------------------
def _grab(pattern, text, cast=float):
    m = re.search(pattern, text, re.I)
    if m:
        raw = m.group(1).replace(",", "").replace("$", "")
        try:
            return cast(raw)
        except ValueError:
            pass
    return None

def parse_int(text, label):
    return _grab(fr"{label}\s*[:|]\s*([\d,]+)", text, int)

def parse_money(text, label):
    return _grab(fr"{label}\s*[:|]\s*\$?([0-9,]+\.\d+)", text, float)

def parse_percent(text, label):
    return _grab(fr"{label}\s*[:|]\s*([\d\.]+)%", text, float)

# ----------------------------------------------------------------------------
# Main PPTX extractor
# ----------------------------------------------------------------------------
def extract_text_from_pptx(file_obj):
    prs = Presentation(io.BytesIO(file_obj.getvalue()))
    structured = []

    kpis = {}  # dict we’ll fill slide‑by‑slide

    for idx, slide in enumerate(prs.slides, 1):
        raw = "".join(extract_text_from_shape(s) for s in slide.shapes)
        stype = identify_slide_type(raw)

        # ---------- Channel‑specific parsing ----------
        if stype == "PMAX_VLA":
            kpis["pmax_vla_impr"] = parse_int(raw,  "Impressions")
            kpis["pmax_vla_clicks"] = parse_int(raw,"Clicks")
            kpis["pmax_vla_cpc"] = parse_money(raw,"CPC")
            kpis["pmax_vla_conv"] = parse_int(raw, "Conversions")
            kpis["pmax_vla_cost_conv"] = parse_money(raw,"Cost / Conversion")

        elif stype == "PMAX":
            kpis["pmax_impr"] = parse_int(raw,  "Impressions")
            kpis["pmax_clicks"] = parse_int(raw,"Clicks")
            kpis["pmax_cpc"] = parse_money(raw,"CPC")
            kpis["pmax_conv"] = parse_int(raw, "Conversions")
            kpis["pmax_cost_conv"] = parse_money(raw,"Cost / Conversion")

        elif stype == "SOCIAL":
            kpis["social_reach"]  = parse_int(raw,  "Reach")
            kpis["social_impr"]   = parse_int(raw,  "Impressions")
            kpis["social_clicks"] = parse_int(raw,  "Clicks")
            kpis["social_cpc"]    = parse_money(raw,"CPC")
            kpis["social_vdp"]    = parse_int(raw,  "VDP Views")

        elif stype == "VIDEO":
            kpis["dv_views"]     = parse_int(raw,  "Views")
            kpis["dv_viewrate"]  = parse_percent(raw,"View Rate")
            kpis["dv_cpc"]       = parse_money(raw,"CPC")
            kpis["dv_cpm"]       = parse_money(raw,"CPM")

        elif stype == "DEMAND_GEN":        # <<< NEW block
            kpis["dg_impr"]   = parse_int  (raw, "Impressions")
            kpis["dg_clicks"] = parse_int  (raw, "Clicks")
            kpis["dg_cpm"]    = parse_money(raw, "CPM")
            kpis["dg_conv"]   = parse_int  (raw, "Conversions")

        elif stype == "BCDF":
            kpis["has_bcdf"] = True
            kpis["bcdf_tactics"] = _grab(r"^(.+?)$", raw.splitlines()[0], str) or ""
            kpis["bcdf_impr"]   = parse_int  (raw, "Impressions")
            kpis["bcdf_clicks"] = parse_int  (raw, "Clicks")
            kpis["bcdf_cpc"]    = parse_money(raw, "CPC")
            kpis["bcdf_vdp"]    = parse_int  (raw, "VDP Views")
            kpis["bcdf_conv"]   = parse_int  (raw, "Conversions")

        # ---------- Write structured dump (for AI path) ----------
        structured.append(f"--- SLIDE {idx} | TYPE: {stype} ---\n{raw}\n" + "-"*80)

    return "\n\n".join(structured), kpis
